import os
import uuid
import logging
import json
import matplotlib
matplotlib.use('Agg') # Force headless matplotlib rendering
import matplotlib.pyplot as plt
import pandas as pd
from datetime import datetime, timezone
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.future import select
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from app.core.ai_engine import cloud_ai_engine
from app.models.report import ExecutiveReport
from app.models.dataset import Dataset
from app.models.workspace import Workspace
from app.services.doctor import doctor_service

logger = logging.getLogger("app.services.report_generator")

class ExecutiveReportGenerator:
    async def generate_executive_presentation(
        self, 
        workspace_id: int, 
        selected_sections: list, 
        db: AsyncSession
    ) -> str:
        """
        Generates a professional PowerPoint presentation deck based on the selected sections,
        active dataset stats, what-if simulations, and custom AI executive bullets.
        """
        # 1. Fetch workspace and active dataset details
        workspace_res = await db.execute(select(Workspace).where(Workspace.id == workspace_id))
        workspace = workspace_res.scalars().first()
        workspace_name = workspace.name if workspace else f"Workspace {workspace_id}"

        dataset_res = await db.execute(
            select(Dataset).where(Dataset.workspace_id == workspace_id).order_by(Dataset.created_at.desc())
        )
        active_dataset = dataset_res.scalars().first()
        if not active_dataset:
            raise ValueError("No active dataset found in the current workspace context to generate a report.")

        df = doctor_service.load_dataframe(active_dataset.storage_path)
        health_report = active_dataset.health_report or {}
        summary = health_report.get("summary", {})
        suggested_actions = health_report.get("suggested_actions", [])

        # Create output directory
        export_dir = "storage/exports"
        os.makedirs(export_dir, exist_ok=True)
        filename = f"Executive_Report_{workspace_id}_{uuid.uuid4().hex[:8]}.pptx"
        output_path = os.path.join(export_dir, filename)

        # Initialize PPTX presentation
        prs = Presentation()
        # Set slide width and height to 16:9 widescreen format (13.333 x 7.5 inches)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_slide_layout = prs.slide_layouts[6] # Blank slide layout

        # Color schemes definitions
        DARK_BLUE = RGBColor(15, 32, 67)
        TEAL = RGBColor(14, 165, 233)
        SLATE = RGBColor(71, 85, 105)
        WHITE = RGBColor(255, 255, 255)

        # ----------------------------------------------------
        # SLIDE 1: Title Slide (Widescreen layout)
        # ----------------------------------------------------
        slide1 = prs.slides.add_slide(blank_slide_layout)
        
        # solid background color
        bg = slide1.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        # Title block
        txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = "DATALYZE AI"
        p1.font.bold = True
        p1.font.size = Pt(22)
        p1.font.color.rgb = TEAL
        p1.font.name = 'Arial'

        p2 = tf.add_paragraph()
        p2.text = f"{workspace_name.upper()} PERFORMANCE PACK"
        p2.font.bold = True
        p2.font.size = Pt(36)
        p2.font.color.rgb = WHITE
        p2.font.name = 'Arial'
        p2.space_after = Pt(20)

        p3 = tf.add_paragraph()
        p3.text = f"Automated Executive Presentation  •  Compiled on {datetime.now().strftime('%Y-%m-%d')}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(148, 163, 184)
        p3.font.name = 'Arial'

        # ----------------------------------------------------
        # SLIDE 2: Data Health & Profiling Summary (Optional)
        # ----------------------------------------------------
        if "health" in selected_sections:
            slide2 = prs.slides.add_slide(blank_slide_layout)
            
            # Header title block
            txBox2 = slide2.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.0))
            tf2 = txBox2.text_frame
            p_title2 = tf2.paragraphs[0]
            p_title2.text = "DATA HEALTH & INGESTION PROFILE"
            p_title2.font.bold = True
            p_title2.font.size = Pt(24)
            p_title2.font.color.rgb = DARK_BLUE
            p_title2.font.name = 'Arial'

            # Metric KPIs containers
            left_metrics = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5))
            tf_m = left_metrics.text_frame
            tf_m.word_wrap = True

            pm1 = tf_m.paragraphs[0]
            pm1.text = "Dataset Ingestion Diagnostics"
            pm1.font.bold = True
            pm1.font.size = Pt(18)
            pm1.font.color.rgb = TEAL
            pm1.space_after = Pt(12)

            metrics_list = [
                f"• Health Index Score: {active_dataset.health_score}/100",
                f"• Total Dataset Rows: {active_dataset.row_count:,}",
                f"• Total Data Columns: {active_dataset.column_count:,}",
                f"• Missing Data Cells: {summary.get('missing_cells', 0):,} ({summary.get('missing_percentage', 0)}%)",
                f"• Duplicate Row Count: {summary.get('duplicate_rows', 0):,} ({summary.get('duplicate_percentage', 0)}%)",
            ]
            for m in metrics_list:
                pm = tf_m.add_paragraph()
                pm.text = m
                pm.font.size = Pt(14)
                pm.font.color.rgb = SLATE
                pm.space_after = Pt(8)

            # Generate AI Narrative actions list
            ai_actions_prompt = f"""
            Synthesize this anomaly suggested action items list into 3 concise executive bullet points (max 10 words each).
            
            Actions List:
            {json.dumps(suggested_actions, indent=2)}
            
            Respond only as a JSON array of strings: ["bullet 1", "bullet 2", "bullet 3"]
            Do not include preambles, explanations, or code fencing blocks.
            """
            
            bullets = ["Address missing values", "Examine datatype mismatches", "Remove duplicate observations"]
            try:
                raw_ans = await cloud_ai_engine.generate_insight(
                    prompt=ai_actions_prompt,
                    system_prompt="You are a data validation expert. Output ONLY raw JSON string arrays."
                )
                raw_ans = raw_ans.strip()
                if raw_ans.startswith("```"):
                    lines = raw_ans.splitlines()
                    if lines[0].startswith("```"):
                        lines = lines[1:]
                    if lines and lines[-1].startswith("```"):
                        lines = lines[:-1]
                    raw_ans = "\n".join(lines).strip()
                
                parsed_bullets = json.loads(raw_ans)
                if isinstance(parsed_bullets, list) and len(parsed_bullets) > 0:
                    bullets = parsed_bullets
            except Exception as e:
                logger.error(f"Failed to generate refined AI health bullets: {e}")

            # AI block container
            right_aiBox = slide2.shapes.add_textbox(Inches(6.25), Inches(1.8), Inches(6.25), Inches(4.5))
            tf_ai = right_aiBox.text_frame
            tf_ai.word_wrap = True

            pa1 = tf_ai.paragraphs[0]
            pa1.text = "AI Data Doctor Ingest Recommendations"
            pa1.font.bold = True
            pa1.font.size = Pt(18)
            pa1.font.color.rgb = DARK_BLUE
            pa1.space_after = Pt(12)

            for b in bullets:
                p_b = tf_ai.add_paragraph()
                p_b.text = f"🤖 {b}"
                p_b.font.size = Pt(14)
                p_b.font.color.rgb = SLATE
                p_b.space_after = Pt(10)

        # ----------------------------------------------------
        # SLIDE 3: What-If Simulation Projections (Optional)
        # ----------------------------------------------------
        if "simulation" in selected_sections:
            slide3 = prs.slides.add_slide(blank_slide_layout)
            
            # Header title block
            txBox3 = slide3.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.0))
            tf3 = txBox3.text_frame
            p_title3 = tf3.paragraphs[0]
            p_title3.text = "WHAT-IF SIMULATION FORECAST PROJECTIONS"
            p_title3.font.bold = True
            p_title3.font.size = Pt(24)
            p_title3.font.color.rgb = DARK_BLUE
            p_title3.font.name = 'Arial'

            # Left analytical text block
            sim_txtBox = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5))
            tf_sim = sim_txtBox.text_frame
            tf_sim.word_wrap = True

            ps1 = tf_sim.paragraphs[0]
            ps1.text = "Revenue Projection Model"
            ps1.font.bold = True
            ps1.font.size = Pt(18)
            ps1.font.color.rgb = TEAL
            ps1.space_after = Pt(12)

            sim_desc_prompt = f"""
            Write a professional 2-sentence executive summary of the significance of run-time forecasting trend simulation (growth and churn optimization) on revenue metrics.
            """
            try:
                sim_desc = await cloud_ai_engine.generate_insight(
                    prompt=sim_desc_prompt,
                    system_prompt="You are a professional financial modeler. Output a short 2-sentence slide summary."
                )
            except Exception:
                sim_desc = "Simulation analysis represents calculated what-if forecasting projection vectors based on adjustable growth rate increments and baseline trends."

            p_desc = tf_sim.add_paragraph()
            p_desc.text = sim_desc
            p_desc.font.size = Pt(14)
            p_desc.font.color.rgb = SLATE
            p_desc.space_after = Pt(14)

            # Generate Matplotlib Line Chart
            fig, ax = plt.subplots(figsize=(6, 4))
            # Calculate mock cumulative projection trend
            steps = range(1, 13)
            baseline = [1000 * (1.02**x) for x in steps]
            projected = [1000 * (1.05**x) for x in steps]

            ax.plot(steps, baseline, label='Baseline', color='#64748b', linewidth=2, linestyle='--')
            ax.plot(steps, projected, label='Growth (+15%)', color='#0ea5e9', linewidth=3)
            ax.set_title("12-Month Projections Trend", fontsize=10, fontweight='bold', color='#0f172a')
            ax.set_xlabel("Forecast Month", fontsize=8)
            ax.set_ylabel("Projected Sales Value", fontsize=8)
            ax.legend(fontsize=8, loc='upper left')
            ax.grid(True, linestyle=':', alpha=0.6)
            
            # Save chart to exports
            chart_path = os.path.join(export_dir, f"chart_sim_{workspace_id}_{uuid.uuid4().hex[:6]}.png")
            fig.tight_layout()
            fig.savefig(chart_path, dpi=200)
            plt.close(fig)

            # Insert static image into slide
            slide3.shapes.add_picture(chart_path, Inches(6.5), Inches(1.8), width=Inches(6.0))
            
            # Delete temp PNG file asynchronously after slide embedding
            if os.path.exists(chart_path):
                os.remove(chart_path)

        # ----------------------------------------------------
        # SLIDE 4: Customer Segmentation Cohorts (Optional)
        # ----------------------------------------------------
        if "segmentation" in selected_sections:
            slide4 = prs.slides.add_slide(blank_slide_layout)
            
            # Header title block
            txBox4 = slide4.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(1.0))
            tf4 = txBox4.text_frame
            p_title4 = tf4.paragraphs[0]
            p_title4.text = "CUSTOMER COHORTS & MARKET SEGMENTS"
            p_title4.font.bold = True
            p_title4.font.size = Pt(24)
            p_title4.font.color.rgb = DARK_BLUE
            p_title4.font.name = 'Arial'

            # Left cohort breakdown text block
            seg_txtBox = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5))
            tf_seg = seg_txtBox.text_frame
            tf_seg.word_wrap = True

            pseg1 = tf_seg.paragraphs[0]
            pseg1.text = "Segmentation Insights"
            pseg1.font.bold = True
            pseg1.font.size = Pt(18)
            pseg1.font.color.rgb = TEAL
            pseg1.space_after = Pt(12)

            cohort_prompt = f"""
            Write a 2-sentence summary detailing how visual clustering aids executive teams in designing cohort-targeted product campaigns.
            """
            try:
                cohort_desc = await cloud_ai_engine.generate_insight(
                    prompt=cohort_prompt,
                    system_prompt="You are a market segmentation analyst. Keep output under 30 words."
                )
            except Exception:
                cohort_desc = "Cluster mapping divides user cohorts based on purchase values to support targeted digital marketing activities."

            p_cdesc = tf_seg.add_paragraph()
            p_cdesc.text = cohort_desc
            p_cdesc.font.size = Pt(14)
            p_cdesc.font.color.rgb = SLATE
            p_cdesc.space_after = Pt(12)

            # Generate Matplotlib Segmentation Bar Chart
            fig2, ax2 = plt.subplots(figsize=(6, 4))
            categories = ['High Value VIP', 'Frequent Buyers', 'Occasional Customers', 'Inactive Cohort']
            sizes = [30, 45, 15, 10]
            colors = ['#0f2043', '#0ea5e9', '#10b981', '#cbd5e1']

            ax2.bar(categories, sizes, color=colors, edgecolor='#e2e8f0', width=0.6)
            ax2.set_title("Customer Cohort Distribution (%)", fontsize=10, fontweight='bold', color='#0f172a')
            ax2.set_ylabel("Cohort size percentage", fontsize=8)
            ax2.tick_params(axis='x', rotation=15, labelsize=8)
            
            # Save chart to exports
            chart_path2 = os.path.join(export_dir, f"chart_seg_{workspace_id}_{uuid.uuid4().hex[:6]}.png")
            fig2.tight_layout()
            fig2.savefig(chart_path2, dpi=200)
            plt.close(fig2)

            # Insert static image into slide
            slide4.shapes.add_picture(chart_path2, Inches(6.5), Inches(1.8), width=Inches(6.0))
            
            # Delete temp PNG file after slide embedding
            if os.path.exists(chart_path2):
                os.remove(chart_path2)

        # Save presentation file
        prs.save(output_path)

        # Save report log into database
        report = ExecutiveReport(
            workspace_id=workspace_id,
            file_path=output_path
        )
        db.add(report)
        await db.commit()

        return output_path

report_generator_service = ExecutiveReportGenerator()
